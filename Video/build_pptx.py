# -*- coding: utf-8 -*-
"""Genera la presentacion PechaKucha (20 slides x 20 s) a partir del guion."""
import struct, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
A = lambda *p: os.path.join(BASE, "assets", *p)

# ---------- paleta ----------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x22, 0x2B, 0x36)
GREY  = RGBColor(0x6B, 0x74, 0x80)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
PARTS = {
    1: {"name": "Integrante 1", "color": RGBColor(0x2E, 0x86, 0xC1)},
    2: {"name": "Integrante 2", "color": RGBColor(0x21, 0x9A, 0x52)},
    3: {"name": "Integrante 3", "color": RGBColor(0xCD, 0x8A, 0x00)},
    4: {"name": "Integrante 4", "color": RGBColor(0xC0, 0x39, 0x2B)},
}
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def set_advance(slide, ms=20000):
    sld = slide._element
    trans = sld.makeelement(qn("p:transition"), {"advClick": "1", "advTm": str(ms)})
    sld.append(trans)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
         first=False, space_after=4, bullet=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Segoe UI"
    return p


def add_fit(slide, path, l, t, w, h, shadow=True):
    iw, ih = png_size(path)
    ar, bar = iw / ih, w / h
    if ar > bar:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    left, top = l + (w - nw) / 2, t + (h - nh) / 2
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(nw), Inches(nh))
    if shadow:
        pic.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        pic.line.width = Pt(0.75)
    return pic


def rbox(slide, l, t, w, h, color, text="", size=13, tcolor=WHITE, bold=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    fill(sp, color)
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for i, ln in enumerate(text.split("\n")):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor
        r.font.name = "Segoe UI"
    return sp


def arrow(slide, x1, y, x2):
    cn = slide.shapes.add_connector(2, Inches(x1), Inches(y), Inches(x2), Inches(y))
    cn.line.color.rgb = GREY; cn.line.width = Pt(2.0)
    le = cn.line._get_or_add_ln()
    he = le.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    le.append(he)


def chrome(slide, idx, part, eje=None, title=""):
    """banda de color + titulo + footer (sin eje ni segundos)."""
    p = PARTS[part]
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(SH))
    fill(band, p["color"])
    # titulo
    tf = textbox(slide, 0.5, 0.34, 11.0, 1.0)
    para(tf, title, 28, p["color"], bold=True, first=True)
    # footer
    ft = textbox(slide, 0.5, SH - 0.45, 12.3, 0.35)
    fp = ft.paragraphs[0]
    fp.alignment = PP_ALIGN.LEFT
    r = fp.add_run(); r.text = f"{p['name']}"
    r.font.size = Pt(10); r.font.color.rgb = GREY; r.font.name = "Segoe UI"
    ft2 = textbox(slide, 0.5, SH - 0.45, 12.3, 0.35)
    fp2 = ft2.paragraphs[0]; fp2.alignment = PP_ALIGN.RIGHT
    r2 = fp2.add_run(); r2.text = f"Diapositiva {idx} / 20"
    r2.font.size = Pt(10); r2.font.color.rgb = GREY; r2.font.name = "Segoe UI"


def bullets(slide, l, t, w, h, items, part, size=15, anchor=MSO_ANCHOR.TOP):
    tf = textbox(slide, l, t, w, h, anchor)
    for i, it in enumerate(items):
        bold = it.startswith("*")
        txt = it.lstrip("*")
        para(tf, txt, size, INK, bold=bold, first=(i == 0), bullet=True, space_after=8)
    return tf


def newslide(idx, part, eje, title):
    s = prs.slides.add_slide(BLANK)
    chrome(s, idx, part, eje, title)
    set_advance(s, 20000)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# =====================================================================
# PARTE 1
# =====================================================================
# --- Slide 1: Portada ---
s = prs.slides.add_slide(BLANK); set_advance(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(1.7))
fill(band, PARTS[1]["color"])
tf = textbox(s, 0.8, 0.35, 11.7, 1.1, MSO_ANCHOR.MIDDLE)
para(tf, "Mantenimiento Predictivo Industrial", 36, WHITE, bold=True, first=True)
para(tf, "Extensión del TP1 — Incorporación de un Perceptrón Multicapa (MLP)", 18, WHITE)
tf2 = textbox(s, 0.8, 2.4, 11.7, 2.6)
para(tf2, "¿Una red neuronal predice mejor la falla que los modelos clásicos?", 22, INK, bold=True, first=True, space_after=18)
para(tf2, "Equipo:  Gaitán · Debona · Agüero · Ocampo", 18, INK, space_after=8)
para(tf2, "Ciencia de Datos / Inteligencia Artificial — UTN FRSF — 2026, 1er Cuatrimestre", 14, GREY)
rbox(s, 0.8, 5.7, 11.7, 0.9, LIGHT,
     "Formato PechaKucha · 20 diapositivas × 20 s = 6:40",
     size=15, tcolor=PARTS[1]["color"])
notes(s, "Buenas. Somos Gaitán, Debona, Agüero y Ocampo. Extendemos nuestro Trabajo "
         "Práctico de mantenimiento predictivo industrial: dado un conjunto de variables "
         "de operación de una máquina, queremos predecir si va a fallar. Hoy sumamos una "
         "red neuronal —un Perceptrón Multicapa— y la comparamos contra los modelos "
         "clásicos que ya entrenamos.")

# --- Slide 2: El dato y el target ---
s = newslide(2, 1, None, "El dato y el target")
add_fit(s, A("cdd", "cdd_02_cell15.png"), 0.55, 1.6, 8.0, 5.1)
bullets(s, 8.7, 1.9, 4.2, 4.6, [
    "*8 variables de entrada",
    "5 continuas: temperatura de aire y de proceso, velocidad, torque y desgaste",
    "Tipo de producto → 3 columnas (One-Hot)",
    "*Target binario: falla / no falla",
    "*Clases balanceadas (~51,5% falla) → sin técnicas de balanceo",
], 1, size=15)
notes(s, "Trabajamos con 8 variables de entrada: cinco continuas —temperatura de aire y de "
         "proceso, velocidad, torque y desgaste de herramienta— y el tipo de producto "
         "codificado en tres columnas. El target es binario: falla o no falla. Las clases "
         "quedaron prácticamente balanceadas, así que no aplicamos técnicas de balanceo.")

# --- Slide 3: Pipeline heredado ---
s = newslide(3, 1, None, "¿Qué heredamos de Ciencia de Datos?")
steps = ["EDA\nexploración", "Limpieza\noutliers · NaN", "One-Hot\nproduct_type",
         "Estandarización\nZ-score", "Train / Test\n80 / 20"]
n = len(steps); bw, gap, y = 2.05, 0.32, 3.0
total = n * bw + (n - 1) * gap
x0 = (SW - total) / 2
for i, st in enumerate(steps):
    x = x0 + i * (bw + gap)
    rbox(s, x, y, bw, 1.25, PARTS[1]["color"] if i % 2 == 0 else RGBColor(0x5B,0xA8,0xD6),
         st, size=13)
    if i < n - 1:
        arrow(s, x + bw + 0.02, y + 0.62, x + bw + gap - 0.02)
rbox(s, x0, 4.9, total, 0.95, LIGHT,
     "Mismos datos para todos los modelos  →  la comparación es JUSTA",
     size=18, tcolor=PARTS[1]["color"])
notes(s, "Reutilizamos el mismo preprocesamiento del TP original: eliminamos outliers, "
         "imputamos faltantes, codificamos product_type con One-Hot —porque demostramos que "
         "es nominal, no ordinal— y estandarizamos con Z-score. Usar exactamente los mismos "
         "datos para todos los modelos es lo que hace justa la comparación que veremos.")

# --- Slide 4: Por qué un MLP ---
s = newslide(4, 1, "Eje 1", "¿Por qué un MLP?")
add_fit(s, A("cdd", "cdd_14_cell36.png"), 0.55, 1.55, 7.2, 5.2)
bullets(s, 8.0, 1.9, 4.9, 4.8, [
    "*Aprende relaciones NO lineales entre variables",
    "El EDA reveló interacciones complejas:",
    "torque bajo dispara la falla",
    "velocidad alta también",
    "y el efecto depende del tipo de producto",
    "*El MLP captura esas combinaciones sin indicárselas",
], 1, size=15)
notes(s, "¿Por qué un Perceptrón Multicapa? Porque, a diferencia de los modelos lineales, "
         "aprende relaciones no lineales entre variables, y en el EDA vimos interacciones "
         "complejas: el torque bajo dispara la falla, la velocidad alta también, y el efecto "
         "depende del tipo de producto. El MLP puede capturar esas combinaciones sin que se "
         "las indiquemos.")

# --- Slide 5: Filosofía de diseño (embudo) ---
s = newslide(5, 1, "Eje 1", "Filosofía de diseño")
widths = [(8, "8 features\nentrada"), (5.0, "32 neuronas"), (3.4, "32 neuronas"), (1.6, "1 salida\nfalla")]
y = 1.9
for i, (wd, lab) in enumerate(widths):
    w_in = 1.2 + wd / 8 * 4.0
    rbox(s, (SW - w_in) / 2, y, w_in, 0.95, PARTS[1]["color"], lab, size=14)
    y += 1.18
bullets(s, 8.7, 2.0, 4.2, 4.4, [
    "*Simplicidad acorde al problema",
    "Con solo 8 features, una red profunda solo agrega sobreajuste",
    "*Arquitectura en embudo: comprime la información progresivamente",
    "Se queda con los patrones que anticipan la falla",
], 1, size=15)
notes(s, "Diseñamos con una idea clave: simplicidad acorde al problema. Con solo 8 features, "
         "una red muy profunda solo agregaría sobreajuste. Optamos por una arquitectura en "
         "embudo: la red comprime progresivamente la información, quedándose con los patrones "
         "que anticipan la falla. Lo concreto lo explica mi compañero.")

# =====================================================================
# PARTE 2
# =====================================================================
# --- Slide 6: Arquitectura final ---
s = newslide(6, 2, "Eje 1", "Arquitectura final")
layers = [("Entrada\n8 neuronas", PARTS[2]["color"]),
          ("Oculta\n32 · ReLU", RGBColor(0x2E,0x7D,0x32)),
          ("Oculta\n32 · ReLU", RGBColor(0x2E,0x7D,0x32)),
          ("Salida\n1 · Sigmoid", RGBColor(0xC0,0x39,0x2B))]
bw, gap, y = 2.45, 0.55, 2.6
total = len(layers) * bw + (len(layers) - 1) * gap
x0 = (SW - total) / 2
for i, (lab, col) in enumerate(layers):
    x = x0 + i * (bw + gap)
    rbox(s, x, y, bw, 1.5, col, lab, size=15)
    if i < len(layers) - 1:
        arrow(s, x + bw + 0.04, y + 0.75, x + bw + gap - 0.04)
rbox(s, x0, 4.7, total, 1.0, LIGHT,
     "Solo 1.377 parámetros  ·  modelo muy liviano  ·  2 capas alcanzan para 8 variables",
     size=17, tcolor=PARTS[2]["color"])
notes(s, "La arquitectura: capa de entrada de 8 neuronas, una por feature; dos capas ocultas "
         "de 32 neuronas; y una única neurona de salida. Dos capas alcanzan para 8 variables; "
         "una sola salida basta porque es clasificación binaria. En total: apenas 1.377 "
         "parámetros, un modelo muy liviano.")

# --- Slide 7: Funciones de activación ---
s = newslide(7, 2, "Eje 1", "Funciones de activación")
add_fit(s, A("gen_activations.png"), 0.55, 1.7, 12.2, 3.6)
bullets(s, 1.0, 5.45, 11.3, 1.6, [
    "*ReLU (capas ocultas): estándar moderno, evita el desvanecimiento del gradiente, muy eficiente",
    "*Sigmoid (salida): mapea a 0–1 → se interpreta como probabilidad de falla.  No usamos Softmax (es para múltiples clases)",
], 2, size=14)
notes(s, "En las capas ocultas usamos ReLU: es el estándar moderno, evita el desvanecimiento "
         "del gradiente y es muy eficiente. En la salida usamos Sigmoid, que mapea el resultado "
         "al rango cero–uno, interpretándolo directamente como probabilidad de falla. No usamos "
         "Softmax porque eso aplica cuando hay varias clases de salida.")

# --- Slide 8: Pérdida y regularización ---
s = newslide(8, 2, "Eje 1", "Función de pérdida y regularización")
rbox(s, 0.7, 2.0, 6.0, 1.4, PARTS[2]["color"],
     "Pérdida:\nBinary Cross-Entropy", size=20)
bullets(s, 0.8, 3.7, 6.0, 2.6, [
    "Función canónica de clasificación binaria",
    "*Penaliza fuerte cuando el modelo está seguro y se equivoca",
], 2, size=15)
rbox(s, 7.0, 2.0, 5.6, 1.4, RGBColor(0x2E,0x7D,0x32),
     "Regularización:\nDropout  +  L2", size=20)
bullets(s, 7.1, 3.7, 5.4, 2.6, [
    "Contempladas para evitar sobreajuste",
    "*La búsqueda automática decide cuánta regularización hace falta realmente",
], 2, size=15)
notes(s, "Como pérdida elegimos binary cross-entropy, la función canónica para clasificación "
         "binaria: penaliza fuerte cuando el modelo está seguro y se equivoca. Para evitar "
         "sobreajuste contemplamos Dropout y regularización L2, dejando que la búsqueda "
         "automática decida cuánta regularización realmente hacía falta.")

# --- Slide 9: Búsqueda de hiperparámetros ---
s = newslide(9, 2, "Eje 1", "Búsqueda de hiperparámetros")
rows = [("Hiperparámetro", "Espacio de búsqueda"),
        ("Neuronas / capa", "16 · 32 · 64"),
        ("Learning rate", "0.001 · 0.01"),
        ("Dropout", "0 · 0.2 · 0.5"),
        ("Regularización L2", "0 · 1e-4 · 1e-3"),
        ("Optimizador", "Adam · RMSprop")]
tbl = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(1.75),
                         Inches(6.3), Inches(3.9)).table
tbl.columns[0].width = Inches(3.0); tbl.columns[1].width = Inches(3.3)
for r, (a, b) in enumerate(rows):
    for c, val in enumerate((a, b)):
        cell = tbl.cell(r, c)
        cell.text = val
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        f = pr.runs[0].font; f.size = Pt(14); f.name = "Segoe UI"
        if r == 0:
            f.bold = True; f.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = PARTS[2]["color"]
        else:
            f.color.rgb = INK
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
bullets(s, 7.3, 2.0, 5.4, 4.6, [
    "*Keras Tuner — RandomSearch",
    "50 ensayos sobre ~1000 combinaciones",
    "*Optimizamos val_loss, NO accuracy",
    "En fallas industriales importa estimar bien la probabilidad, no solo acertar la etiqueta",
], 2, size=15)
notes(s, "No elegimos los hiperparámetros a mano. Usamos Keras Tuner con RandomSearch: 50 "
         "ensayos sobre un espacio de casi mil combinaciones —neuronas, learning rate, dropout, "
         "L2 y optimizador—. Optimizamos sobre val_loss, no accuracy, porque en fallas "
         "industriales nos importa estimar bien la probabilidad, no solo acertar la etiqueta.")

# --- Slide 10: Config óptima y entrenamiento ---
s = newslide(10, 2, "Eje 1", "Configuración óptima y entrenamiento")
add_fit(s, A("ia", "ia_02_cell22.png"), 0.55, 1.55, 8.0, 4.0)
bullets(s, 8.7, 1.8, 4.2, 5.0, [
    "*Adam · lr 0.01 · batch 64",
    "Sin Dropout ni L2: la red era tan chica que no los necesitó",
    "*EarlyStopping: frena al dejar de mejorar",
    "Curvas que convergen bien; validación algo inestable (dataset acotado)",
], 2, size=15)
notes(s, "El mejor modelo: optimizador Adam, learning rate 0.01, batch size 64, sin dropout ni "
         "L2 —la red era tan chica que no los necesitó—. Entrenamos con EarlyStopping para "
         "frenar al dejar de mejorar. Las curvas convergen bien, con una validación algo "
         "inestable, propia de un dataset acotado.")

# =====================================================================
# PARTE 3
# =====================================================================
# --- Slide 11: Sensible a la escala ---
s = newslide(11, 3, "Eje 2", "¿Por qué la red es sensible a la escala?")
add_fit(s, A("cdd", "cdd_04_cell20.png"), 0.55, 1.55, 8.2, 5.2)
bullets(s, 8.9, 1.9, 4.0, 4.8, [
    "*Escalas muy dispares:",
    "RPM rondan los miles",
    "temperaturas apenas decenas",
    "*La red aprende por descenso de gradiente, muy sensible a la escala",
    "Sin tratamiento, las RPM aplastarían a las demás variables",
], 3, size=15)
notes(s, "Acá está el punto central del eje 2. Las variables crudas tienen escalas dispares: "
         "las RPM rondan los miles, las temperaturas apenas decenas. Una red neuronal aprende "
         "por descenso de gradiente, que es muy sensible a la escala: sin tratamiento, las RPM "
         "dominarían el aprendizaje y aplastarían a las demás variables.")

# --- Slide 12: Impacto de estandarizar ---
s = newslide(12, 3, "Eje 2", "El impacto de estandarizar en la red")
rbox(s, 2.6, 2.2, 8.1, 1.5, PARTS[3]["color"], "z  =  ( x − μ )  /  σ", size=30)
rbox(s, 2.6, 3.85, 3.9, 0.85, LIGHT, "media 0", size=18, tcolor=PARTS[3]["color"])
rbox(s, 6.8, 3.85, 3.9, 0.85, LIGHT, "desvío 1", size=18, tcolor=PARTS[3]["color"])
bullets(s, 1.4, 5.0, 10.5, 1.9, [
    "*El gradiente se mueve en terreno parejo → converge más rápido y estable",
    "Todas las variables compiten en igualdad de condiciones",
    "*Para el MLP no es opcional: es condición para que entrene bien",
], 3, size=16)
notes(s, "Por eso aplicamos estandarización Z-score: media cero, desvío uno. El efecto en la "
         "red es directo: el gradiente se mueve en un terreno parejo, converge más rápido y "
         "estable, y todas las variables compiten en igualdad. Para el MLP esto no es opcional: "
         "es condición para que entrene bien.")

# --- Slide 13: Comparación con la Parte 2 ---
s = newslide(13, 3, "Eje 2", "Comparación con los modelos de la Parte 2")
rows = [("Modelo", "¿Sensible a la escala?", "Por qué"),
        ("MLP (red)", "SÍ ✔", "aprende por distancias/gradiente"),
        ("KNN", "SÍ ✔", "se basa en distancias"),
        ("Decision Tree", "NO ✘", "parte por umbrales"),
        ("Random Forest", "NO ✘", "conjunto de árboles, por umbrales")]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.9),
                         Inches(11.9), Inches(3.4)).table
tbl.columns[0].width = Inches(3.0); tbl.columns[1].width = Inches(3.4); tbl.columns[2].width = Inches(5.5)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c); cell.text = val
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        f = pr.runs[0].font; f.size = Pt(15); f.name = "Segoe UI"
        if r == 0:
            f.bold = True; f.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = PARTS[3]["color"]
        else:
            f.color.rgb = INK; f.bold = (c == 1)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
rbox(s, 0.7, 5.65, 11.9, 0.85, LIGHT,
     "Preprocesamiento CRÍTICO para la red  ·  casi indiferente para el ensamblado",
     size=16, tcolor=PARTS[3]["color"])
notes(s, "Comparado con la Parte 2, el contraste es claro. El MLP y KNN dependen de la escala. "
         "En cambio los árboles —Decision Tree y Random Forest— son invariantes: parten por "
         "umbrales, no por distancias, así que estandarizar no cambia sus resultados. Por eso el "
         "preprocesamiento es crítico para la red y casi indiferente para el ensamblado.")

# --- Slide 14: Resultados del MLP ---
s = newslide(14, 3, "Eje 3", "Resultados del MLP")
add_fit(s, A("ia", "ia_03_cell27.png"), 0.55, 1.6, 5.3, 4.4)
add_fit(s, A("ia", "ia_04_cell29.png"), 5.95, 1.6, 4.4, 4.4)
bullets(s, 10.5, 1.9, 2.5, 4.6, [
    "*Acc 96%",
    "*F1 0.96",
    "*AUC 0.991",
    "Umbral 0.3 (no 0.5)",
    "*Recall 98% de fallas detectadas",
], 3, size=14)
rbox(s, 0.7, 6.2, 11.9, 0.65, LIGHT,
     "En mantenimiento, dejar pasar una falla real es mucho más caro que una falsa alarma",
     size=14, tcolor=PARTS[3]["color"])
notes(s, "El MLP rinde muy bien: 96% de accuracy, F1 de 0.96 y un AUC de 0.991. Ajustamos el "
         "umbral de decisión a 0.3 en lugar de 0.5, priorizando el recall: alcanzamos 98% de "
         "fallas detectadas. En mantenimiento, dejar pasar una falla real es mucho más caro que "
         "una falsa alarma.")

# --- Slide 15: Tabla comparativa de métricas ---
s = newslide(15, 3, "Eje 3", "Tabla comparativa de métricas")
add_fit(s, A("cdd", "cdd_22_cell85.png"), 0.55, 1.6, 7.6, 5.0)
bullets(s, 8.4, 1.9, 4.6, 4.8, [
    "Naive Bayes queda atrás: 0.84",
    "KNN y Árbol (base): 0.95–0.96",
    "*MLP: 0.96 · AUC 0.991",
    "*Random Forest lidera: 0.972 · AUC 0.997",
    "Clave: la red, pese a su complejidad, NO supera al ensamblado",
], 3, size=15)
notes(s, "Puestos en una tabla: Naive Bayes queda atrás con 0.84. KNN y Árbol —los modelos "
         "base— rondan 0.95–0.96. El MLP llega a 0.96 con AUC 0.991. Y el ensamblado, Random "
         "Forest, lidera con 0.972 y AUC 0.997. El detalle clave que retoma mi compañero: la "
         "red, pese a su complejidad, no supera al ensamblado.")

# =====================================================================
# PARTE 4
# =====================================================================
# --- Slide 16: MLP frente al ensamblado ---
s = newslide(16, 4, "Eje 3", "MLP frente al ensamblado")
add_fit(s, A("gen_auc.png"), 0.55, 1.6, 8.4, 5.0)
bullets(s, 9.1, 1.9, 3.9, 4.8, [
    "*Random Forest gana en todas las métricas",
    "accuracy · F1 · AUC",
    "aunque por margen estrecho",
    "*Resultado revelador: el modelo más fácil de explicar iguala o supera a la red neuronal",
], 4, size=15)
notes(s, "Comparemos de frente el MLP con el mejor clásico. Random Forest gana en todas las "
         "métricas —accuracy, F1 y AUC—, aunque por un margen estrecho. Es un resultado "
         "revelador: el modelo más fácil de explicar —un conjunto de árboles— iguala o supera a "
         "la red neuronal en este problema.")

# --- Slide 17: ¿Se justifica la complejidad? ---
s = newslide(17, 4, "Eje 3", "¿Se justifica la complejidad del MLP?")
rbox(s, 0.8, 2.0, 5.7, 0.7, PARTS[4]["color"], "COSTO del MLP", size=17)
bullets(s, 0.95, 2.85, 5.5, 3.2, [
    "Estandarización obligatoria",
    "Ajuste de muchos hiperparámetros",
    "Entrenamiento más delicado",
], 4, size=16)
rbox(s, 6.9, 2.0, 5.7, 0.7, RGBColor(0x21,0x9A,0x52), "GANANCIA frente al RF", size=17)
bullets(s, 7.05, 2.85, 5.5, 3.2, [
    "*≈ Ninguna",
    "No mejora al ensamblado",
], 4, size=16)
rbox(s, 0.8, 5.85, 11.8, 0.85, LIGHT,
     "La complejidad solo se paga cuando da una ventaja clara.  Aquí, para este problema:  NO se justifica",
     size=16, tcolor=PARTS[4]["color"])
notes(s, "Entonces, la pregunta de la consigna: ¿la mayor complejidad del MLP se justifica? "
         "Para este problema, no. Exige estandarización obligatoria, ajuste de muchos "
         "hiperparámetros y un entrenamiento más delicado, y aun así no mejora al ensamblado. "
         "La complejidad solo se paga cuando entrega una ventaja clara, y aquí no la hay.")

# --- Slide 18: Precisión vs interpretabilidad ---
s = newslide(18, 4, "Eje 4", "Precisión vs interpretabilidad")
add_fit(s, A("cdd", "cdd_20_cell81.png"), 0.55, 1.6, 8.1, 4.7)
bullets(s, 8.8, 1.9, 4.1, 4.6, [
    "*MLP = caja negra",
    "predice bien, pero no explica por qué",
    "*Árbol / Random Forest = reglas + importancia de cada variable",
    "En planta, el técnico necesita saber QUÉ condición disparó la alarma, no solo un número",
], 4, size=15)
notes(s, "El otro eje de decisión es la interpretabilidad. El MLP es una caja negra: predice "
         "bien, pero no explica por qué. Un árbol o un Random Forest dan reglas e importancia de "
         "cada variable. En una planta, un técnico necesita saber qué condición disparó la "
         "alarma para actuar, no solo recibir un número.")

# --- Slide 19: Modelo elegido para la planta ---
s = newslide(19, 4, "Eje 4", "Modelo elegido para la planta")
rbox(s, 0.7, 1.7, 7.5, 1.0, RGBColor(0x21,0x9A,0x52),
     "Modelo recomendado:  RANDOM FOREST", size=22)
add_fit(s, A("cdd", "cdd_24_cell88.png"), 0.7, 2.95, 7.4, 3.7)
bullets(s, 8.5, 1.9, 4.4, 4.8, [
    "*Mejor performance y robustez",
    "Acc 0.972 · AUC 0.997",
    "*Importancia de variables que respalda cada predicción",
    "Reentrenar cuesta solo 3,7 s → irrelevante frente a esa ventaja",
    "La red queda como alternativa válida, no óptima aquí",
], 4, size=15)
notes(s, "Por eso, para una implementación real en planta, elegimos Random Forest: mejor "
         "performance, robustez e importancia de variables que respalda cada predicción. El "
         "costo de reentrenar, 3,7 segundos, es irrelevante frente a esa ventaja. La red queda "
         "como alternativa válida, pero no como la opción óptima aquí.")

# --- Slide 20: Conclusión ---
s = prs.slides.add_slide(BLANK); set_advance(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(1.3))
fill(band, PARTS[4]["color"])
tf = textbox(s, 0.8, 0.25, 11.7, 0.85, MSO_ANCHOR.MIDDLE)
para(tf, "Conclusión", 34, WHITE, bold=True, first=True)
ideas = [
    "El MLP es potente y exige un preprocesamiento cuidadoso (estandarización obligatoria)",
    "Más complejidad NO garantizó mejores resultados: no supera al ensamblado",
    "Ganó el criterio: elegimos el modelo que combina precisión e interpretabilidad",
]
y = 1.8
for i, t in enumerate(ideas, 1):
    rbox(s, 0.8, y, 0.7, 0.95, PARTS[4]["color"], str(i), size=24)
    tf = textbox(s, 1.7, y, 10.9, 0.95, MSO_ANCHOR.MIDDLE)
    para(tf, t, 18, INK, first=True)
    y += 1.18
rbox(s, 0.8, 5.6, 11.75, 1.1, LIGHT,
     "Gracias.\nGaitán · Debona · Agüero · Ocampo", size=20, tcolor=PARTS[4]["color"])
notes(s, "En síntesis: el MLP es potente y exige un preprocesamiento cuidadoso, pero más "
         "complejidad no garantizó mejores resultados. El criterio ganó: elegimos el modelo que "
         "combina precisión e interpretabilidad para un entorno industrial real. Gracias "
         "—Gaitán, Debona, Agüero y Ocampo.")

out = os.path.join(BASE, "Defensa-PechaKucha-MLP.pptx")
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
